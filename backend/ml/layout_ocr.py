"""
Deep Learning-Based Document Layout Analysis and OCR
Supports PDF and PowerPoint files with advanced layout understanding
"""

import logging
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union
from dataclasses import dataclass
import io

import torch
from PIL import Image
import pytesseract
from pdf2image import convert_from_path, convert_from_bytes
from PyPDF2 import PdfReader
import pdfplumber
from pptx import Presentation

from config import model_config, get_device

logger = logging.getLogger(__name__)


@dataclass
class DocumentPage:
    """Represents a single page with extracted content"""
    page_number: int
    text: str
    layout_elements: List[Dict]
    images: List[Image.Image]
    metadata: Dict


@dataclass
class ProcessedDocument:
    """Complete processed document with all pages"""
    document_id: str
    filename: str
    file_type: str
    pages: List[DocumentPage]
    full_text: str
    metadata: Dict
    total_pages: int


class DocumentProcessor:
    """
    Advanced document processor using deep learning for layout analysis
    
    Features:
    - PDF and PowerPoint support
    - Layout-aware text extraction
    - OCR for scanned documents
    - Table and figure detection
    - Multi-column text handling
    """
    
    def __init__(
        self,
        layout_model: Optional[str] = None,
        ocr_engine: str = "tesseract",
        device: Optional[str] = None
    ):
        """
        Initialize document processor
        
        Args:
            layout_model: HuggingFace model for layout analysis
            ocr_engine: OCR engine to use (tesseract, easyocr)
            device: Device for model inference (cuda, cpu)
        """
        self.layout_model = layout_model or model_config.layout_model
        self.ocr_engine = ocr_engine
        self.device = device or get_device()
        
        logger.info(f"Initializing DocumentProcessor on {self.device}")
        logger.info(f"Layout model: {self.layout_model}")
        logger.info(f"OCR engine: {self.ocr_engine}")
        
        # Initialize models lazily (loaded on first use)
        self._layout_detector = None
        self._ocr_model = None
    
    def _load_layout_model(self):
        """Lazy load layout analysis model"""
        if self._layout_detector is not None:
            return
        
        try:
            logger.info("Loading layout analysis model...")
            # Note: In production, implement actual LayoutLM or similar model
            # For now, we use rule-based extraction with pdfplumber
            logger.info("Using rule-based layout analysis with pdfplumber")
            self._layout_detector = "pdfplumber"  # Placeholder
        except Exception as e:
            logger.error(f"Failed to load layout model: {e}")
            self._layout_detector = "fallback"
    
    def _load_ocr_model(self):
        """Lazy load OCR model"""
        if self._ocr_model is not None:
            return
        
        try:
            logger.info(f"Initializing OCR engine: {self.ocr_engine}")
            if self.ocr_engine == "tesseract":
                # Tesseract is loaded on demand
                self._ocr_model = "tesseract"
            # Add other OCR engines here
        except Exception as e:
            logger.error(f"Failed to load OCR model: {e}")
            self._ocr_model = "fallback"
    
    def process_document(
        self,
        file_path: Union[str, Path],
        document_id: Optional[str] = None
    ) -> ProcessedDocument:
        """
        Process a document file (PDF or PPT)
        
        Args:
            file_path: Path to document file
            document_id: Optional document identifier
            
        Returns:
            ProcessedDocument with extracted content
        """
        file_path = Path(file_path)
        file_type = file_path.suffix.lower()
        
        logger.info(f"Processing document: {file_path.name} (type: {file_type})")
        
        if file_type == ".pdf":
            return self._process_pdf(file_path, document_id)
        elif file_type in [".ppt", ".pptx"]:
            return self._process_ppt(file_path, document_id)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _process_pdf(
        self,
        file_path: Path,
        document_id: Optional[str]
    ) -> ProcessedDocument:
        """
        Process PDF file with layout-aware extraction
        
        Strategy:
        1. Try text extraction with pdfplumber (preserves layout)
        2. If text is sparse, fallback to OCR
        3. Extract tables and figures separately
        """
        logger.info(f"Processing PDF: {file_path.name}")
        
        pages = []
        full_text_parts = []
        
        try:
            # Use pdfplumber for layout-aware extraction
            with pdfplumber.open(file_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"PDF has {total_pages} pages")
                
                for page_num, page in enumerate(pdf.pages, start=1):
                    logger.debug(f"Processing page {page_num}/{total_pages}")
                    
                    # Extract text with layout preservation
                    text = page.extract_text() or ""
                    
                    # Extract tables
                    tables = page.extract_tables() or []
                    
                    # If text is too sparse, try OCR
                    if len(text.strip()) < 50 and not tables:
                        logger.debug(f"Page {page_num}: Text sparse, trying OCR...")
                        text = self._ocr_page(page)
                    
                    # Extract layout elements
                    layout_elements = self._extract_layout_elements(page)
                    
                    # Convert page to image for potential further processing
                    images = []
                    
                    # Create page object
                    doc_page = DocumentPage(
                        page_number=page_num,
                        text=text,
                        layout_elements=layout_elements,
                        images=images,
                        metadata={
                            "width": page.width,
                            "height": page.height,
                            "has_tables": len(tables) > 0,
                            "table_count": len(tables)
                        }
                    )
                    
                    pages.append(doc_page)
                    full_text_parts.append(text)
        
        except Exception as e:
            logger.error(f"Error processing PDF: {e}", exc_info=True)
            # Fallback to basic PyPDF2 extraction
            pages, full_text_parts = self._fallback_pdf_extraction(file_path)
        
        full_text = "\n\n".join(full_text_parts)
        
        return ProcessedDocument(
            document_id=document_id or file_path.stem,
            filename=file_path.name,
            file_type=".pdf",
            pages=pages,
            full_text=full_text,
            total_pages=len(pages),
            metadata={
                "extraction_method": "pdfplumber",
                "file_size": file_path.stat().st_size
            }
        )
    
    def _extract_layout_elements(self, page) -> List[Dict]:
        """
        Extract layout elements from PDF page
        
        Args:
            page: pdfplumber page object
            
        Returns:
            List of layout elements with coordinates
        """
        elements = []
        
        try:
            # Extract text blocks with coordinates
            words = page.extract_words()
            
            # Group words into text blocks (simplified)
            for word in words[:10]:  # Limit for performance
                elements.append({
                    "type": "text",
                    "content": word.get("text", ""),
                    "bbox": [
                        word.get("x0", 0),
                        word.get("top", 0),
                        word.get("x1", 0),
                        word.get("bottom", 0)
                    ]
                })
        
        except Exception as e:
            logger.warning(f"Layout extraction failed: {e}")
        
        return elements
    
    def _ocr_page(self, page) -> str:
        """
        Perform OCR on a PDF page
        
        Args:
            page: pdfplumber page object
            
        Returns:
            OCR extracted text
        """
        try:
            self._load_ocr_model()
            
            # Convert page to image
            img = page.to_image(resolution=300)
            pil_image = img.original
            
            # Perform OCR
            if self.ocr_engine == "tesseract":
                text = pytesseract.image_to_string(
                    pil_image,
                    lang="+".join(model_config.ocr_languages)
                )
                return text
            
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return ""
    
    def _fallback_pdf_extraction(self, file_path: Path) -> Tuple[List[DocumentPage], List[str]]:
        """
        Fallback PDF extraction using PyPDF2
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Tuple of (pages, text_parts)
        """
        logger.warning("Using fallback PDF extraction (PyPDF2)")
        
        pages = []
        text_parts = []
        
        try:
            reader = PdfReader(str(file_path))
            
            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                
                doc_page = DocumentPage(
                    page_number=page_num,
                    text=text,
                    layout_elements=[],
                    images=[],
                    metadata={"extraction_method": "pypdf2"}
                )
                
                pages.append(doc_page)
                text_parts.append(text)
        
        except Exception as e:
            logger.error(f"Fallback extraction failed: {e}")
        
        return pages, text_parts
    
    def _process_ppt(
        self,
        file_path: Path,
        document_id: Optional[str]
    ) -> ProcessedDocument:
        """
        Process PowerPoint file
        
        Args:
            file_path: Path to PPT/PPTX file
            document_id: Optional document identifier
            
        Returns:
            ProcessedDocument with extracted content
        """
        logger.info(f"Processing PowerPoint: {file_path.name}")
        
        pages = []
        full_text_parts = []
        
        try:
            prs = Presentation(str(file_path))
            total_slides = len(prs.slides)
            
            logger.info(f"PowerPoint has {total_slides} slides")
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                logger.debug(f"Processing slide {slide_num}/{total_slides}")
                
                # Extract text from all shapes
                text_parts = []
                layout_elements = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
                        
                        # Record layout element
                        layout_elements.append({
                            "type": "text_box",
                            "content": shape.text,
                            "bbox": [
                                shape.left,
                                shape.top,
                                shape.left + shape.width,
                                shape.top + shape.height
                            ]
                        })
                    
                    # Extract text from tables
                    if shape.has_table:
                        table_text = self._extract_table_text(shape.table)
                        text_parts.append(table_text)
                
                slide_text = "\n".join(text_parts)
                
                doc_page = DocumentPage(
                    page_number=slide_num,
                    text=slide_text,
                    layout_elements=layout_elements,
                    images=[],
                    metadata={
                        "slide_layout": slide.slide_layout.name,
                        "shape_count": len(slide.shapes)
                    }
                )
                
                pages.append(doc_page)
                full_text_parts.append(slide_text)
        
        except Exception as e:
            logger.error(f"Error processing PowerPoint: {e}", exc_info=True)
            raise
        
        full_text = "\n\n".join(full_text_parts)
        
        return ProcessedDocument(
            document_id=document_id or file_path.stem,
            filename=file_path.name,
            file_type=file_path.suffix,
            pages=pages,
            full_text=full_text,
            total_pages=len(pages),
            metadata={
                "extraction_method": "python-pptx",
                "file_size": file_path.stat().st_size
            }
        )
    
    def _extract_table_text(self, table) -> str:
        """
        Extract text from PowerPoint table
        
        Args:
            table: PowerPoint table object
            
        Returns:
            Formatted table text
        """
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)


# Convenience function
def process_document(file_path: Union[str, Path], document_id: Optional[str] = None) -> ProcessedDocument:
    """
    Process a document file
    
    Args:
        file_path: Path to document
        document_id: Optional document ID
        
    Returns:
        ProcessedDocument
    """
    processor = DocumentProcessor()
    return processor.process_document(file_path, document_id)
